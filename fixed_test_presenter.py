import streamlit as st
import pandas as pd
from datetime import datetime

# Test if basic imports work
try:
    from pptx import Presentation
    pptx_available = True
except ImportError as e:
    pptx_available = False

try:
    import pyodbc
    pyodbc_available = True
except ImportError as e:
    pyodbc_available = False

try:
    import openai
    openai_available = True
except ImportError as e:
    openai_available = False

try:
    import plotly.express as px
    plotly_available = True
except ImportError as e:
    plotly_available = False

# Page configuration - REMOVED EMOJI to fix the error
st.set_page_config(
    page_title="Mental Health AI Test",
    layout="wide"
)

st.title("Mental Health AI Platform - Library Test")
st.write(f"Test run at: {datetime.now()}")

# Display library status
st.header("Library Status Check")

status_data = []

# Check each library
if pptx_available:
    st.success("PowerPoint library (python-pptx) is working!")
    status_data.append({"Library": "python-pptx", "Status": "Working"})
else:
    st.error("PowerPoint library (python-pptx) has issues")
    status_data.append({"Library": "python-pptx", "Status": "Error"})

if pyodbc_available:
    st.success("Database library (pyodbc) is working!")
    status_data.append({"Library": "pyodbc", "Status": "Working"})
else:
    st.error("Database library (pyodbc) has issues")
    status_data.append({"Library": "pyodbc", "Status": "Error"})

if openai_available:
    st.success("OpenAI library is working!")
    status_data.append({"Library": "openai", "Status": "Working"})
else:
    st.error("OpenAI library has issues")
    status_data.append({"Library": "openai", "Status": "Error"})

if plotly_available:
    st.success("Plotly library is working!")
    status_data.append({"Library": "plotly", "Status": "Working"})
else:
    st.error("Plotly library has issues")
    status_data.append({"Library": "plotly", "Status": "Error"})

# Always working
status_data.append({"Library": "streamlit", "Status": "Working"})
status_data.append({"Library": "pandas", "Status": "Working"})

# Show status table
st.dataframe(pd.DataFrame(status_data))

# Test basic Streamlit functionality
st.header("Streamlit Functionality Test")

if st.button("Test Button"):
    st.success("Button clicked! Streamlit is working correctly!")

# Test sidebar
with st.sidebar:
    st.header("Test Sidebar")
    st.write("If you can see this, Streamlit sidebar is working!")
    
    api_key = st.text_input("Test OpenAI API Key", type="password")
    if api_key:
        st.success("API key entered successfully!")

# Test PowerPoint creation
if pptx_available:
    st.header("PowerPoint Test")
    if st.button("Test PowerPoint Creation"):
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = "Test Slide from Mental Health AI Platform"
            
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            
            # Save to memory
            from io import BytesIO
            ppt_buffer = BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            st.success("PowerPoint creation test passed!")
            st.download_button(
                label="Download Test PowerPoint",
                data=ppt_buffer.getvalue(),
                file_name="test_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"PowerPoint creation failed: {e}")
else:
    st.header("PowerPoint Test")
    st.warning("PowerPoint library not available - install with: conda install -c conda-forge python-pptx")

# Test plotting
if plotly_available:
    st.header("Visualization Test")
    chart_data = pd.DataFrame({
        'condition': ['Depression', 'Anxiety', 'Bipolar', 'PTSD'],
        'count': [45, 32, 12, 8],
        'avg_sessions': [8.5, 6.2, 12.1, 15.3]
    })
    
    fig = px.bar(chart_data, x='condition', y='count', title='Mental Health Conditions Test Chart')
    st.plotly_chart(fig, use_container_width=True)
    
    # Test second chart
    fig2 = px.line(chart_data, x='condition', y='avg_sessions', title='Average Sessions by Condition')
    st.plotly_chart(fig2, use_container_width=True)
else:
    st.header("Visualization Test")
    st.warning("Plotly library not available - install with: conda install -c conda-forge plotly")

# Test data manipulation
st.header("Data Processing Test")
test_data = pd.DataFrame({
    'patient_id': ['001', '002', '003', '004', '005'],
    'condition': ['Depression', 'Anxiety', 'Depression', 'Bipolar', 'PTSD'],
    'sessions': [8, 6, 12, 15, 10],
    'progress': ['Good', 'Excellent', 'Fair', 'Good', 'Excellent']
})

st.write("Sample Mental Health Data:")
st.dataframe(test_data)

# Test CSV download
csv = test_data.to_csv(index=False)
st.download_button(
    label="Download Test Data (CSV)",
    data=csv,
    file_name="test_mental_health_data.csv",
    mime="text/csv"
)

# Overall status
st.header("Overall System Status")

working_libraries = sum([pptx_available, pyodbc_available, openai_available, plotly_available])
total_libraries = 4

if working_libraries == total_libraries:
    st.success(f"All {total_libraries} libraries are working! You're ready for the full Mental Health AI Platform.")
    st.info("You can now run your main presenter_agent.py file.")
elif working_libraries >= 2:
    st.warning(f"{working_libraries}/{total_libraries} libraries working. Some features may be limited.")
    st.info("Install missing libraries and try again.")
else:
    st.error(f"Only {working_libraries}/{total_libraries} libraries working. Please install missing dependencies.")

# Installation commands
st.header("Installation Commands")
st.write("If any libraries are missing, run these commands:")

st.code("""
# Activate conda environment
conda activate base

# Install missing packages
conda install -c conda-forge python-pptx pyodbc plotly streamlit
pip install openai
""")

# Next steps
st.header("Next Steps")
if working_libraries >= 3:
    st.success("System is ready! Try running your main Mental Health AI application.")
    st.info("Replace this test file with your presenter_agent.py code and run with:")
    st.code("streamlit run presenter_agent.py")
else:
    st.warning("Install missing libraries first, then try running your main application.")

st.write("---")
st.write("Mental Health AI Platform Test - Transforming Lives")
st.write(f"Test completed at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
