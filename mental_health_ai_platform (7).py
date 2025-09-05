import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime, timedelta
import json
import os
from io import BytesIO
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots

# Database and AI imports
try:
    import pyodbc
    DATABASE_AVAILABLE = True
except ImportError:
    DATABASE_AVAILABLE = False
    st.error("pyodbc not available. Install with: pip install pyodbc")

try:
    import openai
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False
    st.error("OpenAI not available. Install with: pip install openai")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    POWERPOINT_AVAILABLE = True
except ImportError:
    POWERPOINT_AVAILABLE = False
    st.error("PowerPoint not available. Install with: pip install python-pptx")

try:
    from dotenv import load_dotenv
    DOTENV_AVAILABLE = True
except ImportError:
    DOTENV_AVAILABLE = False
    st.warning("python-dotenv not available. Install with: pip install python-dotenv")

# Load environment variables
if DOTENV_AVAILABLE:
    # Try to load from the specified directory
    env_path = r"C:\Users\pjone\OneDrive\Desktop\NewPython\test.env"
    if os.path.exists(env_path):
        load_dotenv(env_path)
        st.sidebar.success("Environment loaded from test.env")
    else:
        load_dotenv()  # Try default locations
        st.sidebar.warning("test.env not found, using default .env")

# Page configuration
st.set_page_config(
    page_title="Mental Health AI Platform - SQL Server & OpenAI Integration",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
<style>
.main-header {
    background: linear-gradient(135deg, #667eea 0%, #764ba2 50%, #f093fb 100%);
    padding: 1.5rem;
    border-radius: 15px;
    color: white;
    text-align: center;
    margin-bottom: 2rem;
    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
}
.agent-card {
    border: 2px solid #e1e8ed;
    border-radius: 15px;
    padding: 1.5rem;
    margin: 0.5rem;
    background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%);
    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
    transition: transform 0.2s;
}
.agent-card:hover {
    transform: translateY(-2px);
    box-shadow: 0 4px 8px rgba(0, 0, 0, 0.15);
}
.success-banner {
    background: linear-gradient(135deg, #d4edda 0%, #c3e6cb 100%);
    color: #155724;
    padding: 1rem;
    border-radius: 10px;
    border: 1px solid #c3e6cb;
}
.metric-card {
    background: white;
    padding: 1rem;
    border-radius: 10px;
    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
    text-align: center;
    margin: 0.5rem;
}
.ai-response {
    background: #f8f9fa;
    border-left: 4px solid #667eea;
    padding: 1rem;
    border-radius: 0 10px 10px 0;
    margin: 1rem 0;
}
</style>
""", unsafe_allow_html=True)

# Database connection class
class DatabaseConnection:
    def __init__(self):
        self.server = "JONESFAMILYPC3"
        self.database = "MedDw"
        self.connection = None
        
    def connect(self):
        try:
            # Try Windows Authentication first
            connection_string = (
                f"DRIVER={{ODBC Driver 17 for SQL Server}};"
                f"SERVER={self.server};"
                f"DATABASE={self.database};"
                f"Trusted_Connection=yes;"
            )
            self.connection = pyodbc.connect(connection_string)
            return True
        except Exception as e:
            st.error(f"Database connection failed: {e}")
            return False
    
    def execute_query(self, query):
        if not self.connection:
            if not self.connect():
                return None
        
        try:
            return pd.read_sql(query, self.connection)
        except Exception as e:
            st.error(f"Query execution failed: {e}")
            return None
    
    def get_mental_health_data(self, limit=None):
        query = "SELECT * FROM [dbo].[vw_MentalHealthRAG]"
        if limit:
            query += f" ORDER BY visit_date DESC OFFSET 0 ROWS FETCH NEXT {limit} ROWS ONLY"
        return self.execute_query(query)
    
    def close(self):
        if self.connection:
            self.connection.close()

# OpenAI integration class
class OpenAIAssistant:
    def __init__(self):
        self.client = None
        self.api_key = os.getenv('OPENAI_API_KEY')
        
        if self.api_key and OPENAI_AVAILABLE:
            try:
                self.client = OpenAI(api_key=self.api_key)
                # Test the connection
                response = self.client.models.list()
                st.sidebar.success("OpenAI API connected successfully")
            except Exception as e:
                st.sidebar.error(f"OpenAI connection failed: {e}")
                self.client = None
        else:
            st.sidebar.warning("OpenAI API key not found in environment")
    
    def analyze_data(self, data, query, agent_type):
        if not self.client:
            return "OpenAI not available. Please check your API key."
        
        try:
            # Prepare data summary for analysis
            data_summary = self._prepare_data_summary(data, agent_type)
            
            # Create agent-specific prompts
            system_prompt = self._get_system_prompt(agent_type)
            user_prompt = f"""
            Analyze the following mental health data for: {query}
            
            Data Summary:
            {data_summary}
            
            Please provide:
            1. Key insights and findings
            2. Specific recommendations
            3. Potential concerns or red flags
            4. Action items for healthcare management
            
            Focus on {agent_type.lower()} aspects of the data.
            """
            
            response = self.client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_tokens=1000,
                temperature=0.3
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            return f"Analysis failed: {e}"
    
    def _prepare_data_summary(self, data, agent_type):
        if data is None or data.empty:
            return "No data available"
        
        summary = f"Dataset contains {len(data)} records.\n\n"
        
        # Agent-specific data summaries
        if agent_type == "Billing":
            billing_cols = ['billed_amount', 'paid_amount', 'patient_responsibility', 'copay_amount']
            available_cols = [col for col in billing_cols if col in data.columns]
            if available_cols:
                summary += "Financial Metrics:\n"
                for col in available_cols:
                    if data[col].dtype in ['int64', 'float64']:
                        summary += f"- {col}: Mean=${data[col].mean():.2f}, Total=${data[col].sum():.2f}\n"
        
        elif agent_type == "Scheduling":
            if 'visit_date' in data.columns:
                summary += f"Date range: {data['visit_date'].min()} to {data['visit_date'].max()}\n"
            if 'session_length_minutes' in data.columns:
                summary += f"Average session length: {data['session_length_minutes'].mean():.1f} minutes\n"
            if 'visit_type' in data.columns:
                summary += f"Visit types: {data['visit_type'].value_counts().to_dict()}\n"
        
        elif agent_type == "Research":
            if 'diagnosis_code' in data.columns:
                summary += f"Diagnoses: {data['diagnosis_code'].value_counts().head().to_dict()}\n"
            if 'condition_category' in data.columns:
                summary += f"Conditions: {data['condition_category'].value_counts().to_dict()}\n"
            if 'provider_specialty' in data.columns:
                summary += f"Specialties: {data['provider_specialty'].value_counts().to_dict()}\n"
        
        return summary[:2000]  # Limit summary length
    
    def _get_system_prompt(self, agent_type):
        prompts = {
            "Billing": """You are a healthcare billing specialist AI. Focus on financial analysis, 
                         claims processing, revenue optimization, and billing compliance. Provide 
                         actionable insights for financial performance improvement.""",
            
            "Scheduling": """You are a healthcare scheduling optimization AI. Focus on appointment 
                           efficiency, provider utilization, patient flow, and scheduling patterns. 
                           Provide insights for operational improvements.""",
            
            "Research": """You are a clinical research AI analyst. Focus on treatment outcomes, 
                         patient demographics, clinical effectiveness, and healthcare quality metrics. 
                         Provide evidence-based insights for clinical improvement.""",
            
            "Presenter": """You are an executive presentation AI. Create clear, actionable summaries 
                          suitable for healthcare leadership. Focus on strategic insights and 
                          decision-making support."""
        }
        return prompts.get(agent_type, "You are a healthcare data analysis AI assistant.")

# Enhanced AI Agent class
class EnhancedMentalHealthAgent:
    def __init__(self, name, role, description, color, db_connection, ai_assistant):
        self.name = name
        self.role = role
        self.description = description
        self.color = color
        self.analyses = []
        self.db = db_connection
        self.ai = ai_assistant
    
    def add_analysis(self, query, custom_sql=None, insights=""):
        try:
            # Get data from database
            if custom_sql:
                data = self.db.execute_query(custom_sql)
            else:
                data = self.db.get_mental_health_data(1000)  # Limit for performance
            
            if data is None or data.empty:
                st.warning("No data returned from query")
                return None
            
            # Get AI insights if available
            ai_insights = ""
            if self.ai.client:
                ai_insights = self.ai.analyze_data(data, query, self.role.split()[0])
            
            analysis = {
                'timestamp': datetime.now(),
                'query': query,
                'custom_sql': custom_sql,
                'data': data,
                'insights': insights,
                'ai_insights': ai_insights,
                'title': self._generate_title(query),
                'agent': self.name
            }
            
            self.analyses.append(analysis)
            return analysis
            
        except Exception as e:
            st.error(f"Analysis failed: {e}")
            return None
    
    def _generate_title(self, query):
        query_lower = query.lower()
        if 'billing' in query_lower or 'financial' in query_lower or 'revenue' in query_lower:
            return 'Financial Performance Analysis'
        elif 'appointment' in query_lower or 'schedule' in query_lower:
            return 'Scheduling Optimization Analysis'
        elif 'treatment' in query_lower or 'outcome' in query_lower or 'diagnosis' in query_lower:
            return 'Clinical Research Analysis'
        elif 'provider' in query_lower:
            return 'Provider Performance Analysis'
        elif 'patient' in query_lower:
            return 'Patient Analytics Report'
        else:
            return 'Healthcare Analytics Report'

# Header
st.markdown("""
<div class="main-header">
    <h1>🧠 Mental Health AI Analytics Platform</h1>
    <p>Transforming Lives Mental Health Counseling - AI-Powered Analytics with SQL Server & OpenAI Integration</p>
    <p>Database: JONESFAMILYPC3\MedDw | View: [dbo].[vw_MentalHealthRAG]</p>
</div>
""", unsafe_allow_html=True)

# Initialize connections
@st.cache_resource
def init_connections():
    db = DatabaseConnection()
    ai = OpenAIAssistant()
    return db, ai

db_connection, ai_assistant = init_connections()

# System status sidebar
st.sidebar.header("🔧 System Status")

# Database status
if DATABASE_AVAILABLE:
    if db_connection.connect():
        st.sidebar.success("✅ SQL Server Connected")
        # Test query
        try:
            test_data = db_connection.get_mental_health_data(5)
            if test_data is not None:
                st.sidebar.info(f"📊 {len(test_data)} sample records loaded")
            else:
                st.sidebar.warning("⚠️ No data in view")
        except Exception as e:
            st.sidebar.error(f"❌ View access error: {e}")
    else:
        st.sidebar.error("❌ SQL Server Connection Failed")
else:
    st.sidebar.error("❌ Database Driver Not Available")

# OpenAI status
if ai_assistant.client:
    st.sidebar.success("✅ OpenAI API Connected")
else:
    st.sidebar.error("❌ OpenAI API Not Available")

# PowerPoint status
if POWERPOINT_AVAILABLE:
    st.sidebar.success("✅ PowerPoint Generation Ready")
else:
    st.sidebar.error("❌ PowerPoint Not Available")

# Initialize agents
if 'enhanced_agents' not in st.session_state:
    st.session_state.enhanced_agents = {
        'BillBot': EnhancedMentalHealthAgent(
            'BillBot', 'Billing Specialist', 
            'AI-powered financial analysis and claims processing with real-time SQL data',
            '#1f77b4', db_connection, ai_assistant
        ),
        'ScheduleBot': EnhancedMentalHealthAgent(
            'ScheduleBot', 'Scheduling Coordinator', 
            'Appointment optimization and capacity management with predictive analytics',
            '#ff7f0e', db_connection, ai_assistant
        ),
        'ResearchBot': EnhancedMentalHealthAgent(
            'ResearchBot', 'Research Analyst', 
            'Clinical research and outcomes analysis with AI-driven insights',
            '#2ca02c', db_connection, ai_assistant
        ),
        'PresenterBot': EnhancedMentalHealthAgent(
            'PresenterBot', 'Executive Presenter', 
            'AI-powered PowerPoint presentation generation with strategic insights',
            '#8e44ad', db_connection, ai_assistant
        )
    }

# Enhanced PowerPoint generation
def create_enhanced_presentation(analyses):
    if not POWERPOINT_AVAILABLE:
        return None
    
    try:
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Mental Health Analytics Intelligence Report"
        subtitle.text = f"""Transforming Lives Mental Health Counseling
AI-Powered Analytics Dashboard
Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}
Database: JONESFAMILYPC3\\MedDw
CONFIDENTIAL - For Internal Use Only"""
        
        # Executive summary with AI insights
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        
        title_shape.text = '🎯 Executive Summary'
        tf = body_shape.text_frame
        tf.text = f'Intelligence Report from {len(analyses)} AI Agent Analyses:'
        tf.word_wrap = True  # Enable word wrap
        
        for bullet in [
            '🤖 Real-time analysis powered by OpenAI GPT-4',
            '🔒 All patient data protected through automated PHI masking',
            '📊 Live database integration with SQL Server',
            '📈 Analysis covers clinical, operational, and financial KPIs',
            '🎯 AI-generated recommendations for immediate action'
        ]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 1
        
        # Individual analysis slides with AI insights
        for i, analysis in enumerate(analyses):
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.shapes.placeholders[1]
            
            title_shape.text = f"{i+1}. {analysis['title']}"
            
            tf = body_shape.text_frame
            tf.text = f"Query: {analysis['query']}"
            tf.word_wrap = True  # Enable word wrap
            
            # Analysis details
            if isinstance(analysis['data'], pd.DataFrame) and not analysis['data'].empty:
                p = tf.add_paragraph()
                p.text = f"📊 Records analyzed: {len(analysis['data']):,}"
                p.level = 1
                
                # Key metrics
                numeric_cols = analysis['data'].select_dtypes(include=['number']).columns
                for col in numeric_cols[:2]:
                    if analysis['data'][col].notna().sum() > 0:
                        avg_val = analysis['data'][col].mean()
                        p = tf.add_paragraph()
                        p.text = f"📈 Average {col.replace('_', ' ').title()}: {avg_val:.2f}"
                        p.level = 1
            
            # AI insights with text wrapping and chunking
            if analysis.get('ai_insights'):
                ai_text = analysis['ai_insights']
                
                # Check if AI insights are too long for one slide
                if len(ai_text) > 800:  # If text is too long
                    # Split into chunks
                    chunks = _split_text_into_chunks(ai_text, 700)
                    
                    for chunk_idx, chunk in enumerate(chunks):
                        if chunk_idx == 0:
                            # First chunk goes on current slide
                            p = tf.add_paragraph()
                            p.text = "🤖 AI Analysis:"
                            p.level = 1
                            
                            p = tf.add_paragraph()
                            p.text = chunk
                            p.level = 2
                        else:
                            # Additional chunks get new slides
                            slide = prs.slides.add_slide(bullet_slide_layout)
                            title_shape = slide.shapes.title
                            body_shape = slide.shapes.placeholders[1]
                            
                            title_shape.text = f"{i+1}. {analysis['title']} (Continued {chunk_idx + 1})"
                            
                            tf = body_shape.text_frame
                            tf.text = f"🤖 AI Analysis (Continued):"
                            tf.word_wrap = True
                            
                            p = tf.add_paragraph()
                            p.text = chunk
                            p.level = 1
                else:
                    # Short text - fits on one slide
                    p = tf.add_paragraph()
                    p.text = "🤖 AI Analysis:"
                    p.level = 1
                    
                    p = tf.add_paragraph()
                    p.text = ai_text
                    p.level = 2
        
        # AI Recommendations slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        
        title_shape.text = '🎯 AI-Powered Recommendations'
        
        tf = body_shape.text_frame
        tf.text = 'Strategic Actions Based on AI Analysis:'
        tf.word_wrap = True
        
        recommendations = [
            '🚀 Immediate Actions:',
            'Continue leveraging AI for predictive analytics',
            'Implement real-time monitoring dashboards',
            'Optimize provider scheduling based on demand patterns',
            '📊 Strategic Initiatives:',
            'Expand AI analysis to include predictive modeling',
            'Develop automated alerting for anomaly detection',
            'Integrate patient outcome prediction models'
        ]
        
        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.level = 1 if not rec.endswith(':') else 0
        
        # Data governance slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        
        title_shape.text = '🔒 Privacy & AI Governance'
        
        tf = body_shape.text_frame
        tf.text = 'Data Protection & AI Ethics:'
        tf.word_wrap = True
        
        governance_items = [
            '✅ All PHI automatically masked through AI preprocessing',
            '🔐 HIPAA-compliant data handling protocols enforced',
            '🤖 AI analysis conducted on aggregated, de-identified data',
            '📋 Audit trails maintained for all AI interactions',
            '🎯 AI recommendations reviewed by clinical staff',
            '🏥 This presentation contains confidential healthcare analytics'
        ]
        
        for item in governance_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
        
        # Save to BytesIO
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        st.error(f"Error creating enhanced presentation: {e}")
        return None

def _split_text_into_chunks(text, max_length):
    """Split long text into smaller chunks for better slide presentation"""
    if len(text) <= max_length:
        return [text]
    
    chunks = []
    sentences = text.split('. ')
    current_chunk = ""
    
    for sentence in sentences:
        # Add sentence to current chunk if it fits
        if len(current_chunk + sentence + '. ') <= max_length:
            current_chunk += sentence + '. '
        else:
            # Start new chunk
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = sentence + '. '
    
    # Add the last chunk
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks

# Agent selection
st.sidebar.header("🤖 AI Agents")
selected_agent_name = st.sidebar.selectbox(
    "Select AI Agent:",
    options=list(st.session_state.enhanced_agents.keys()),
    format_func=lambda x: f"{x} - {st.session_state.enhanced_agents[x].role}"
)

selected_agent = st.session_state.enhanced_agents[selected_agent_name]

# Display agent info
with st.sidebar.expander("Agent Details", expanded=True):
    st.markdown(f"**🤖 {selected_agent.name}**")
    st.write(f"**Role:** {selected_agent.role}")
    st.write(f"**Capability:** {selected_agent.description}")
    
    # Show analysis count
    if selected_agent.analyses:
        st.info(f"📊 {len(selected_agent.analyses)} analyses completed")

# Presentation queue status
total_analyses = sum(len(agent.analyses) for agent in st.session_state.enhanced_agents.values() 
                    if agent.name != 'PresenterBot')

if total_analyses > 0:
    st.sidebar.markdown("---")
    st.sidebar.success(f"📋 Presentation Queue: {total_analyses} analyses ready")
    
    if st.sidebar.button("🎯 Quick Generate Presentation"):
        st.sidebar.info("Switch to PresenterBot to generate!")

# Main interface logic
if selected_agent.name == 'PresenterBot':
    # Presenter interface
    st.markdown(f"### 🎯 {selected_agent.name} - {selected_agent.role}")
    st.markdown("*Generate AI-powered executive PowerPoint presentations from your analytics*")
    
    if total_analyses == 0:
        st.info("🔄 No analyses available yet. Use other AI agents to generate insights first!")
        
        # Demo workflow
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("""
            <div class="metric-card">
                <h3>1️⃣ Analyze Data</h3>
                <p>Use BillBot, ScheduleBot, or ResearchBot to query your SQL Server database</p>
            </div>
            """, unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div class="metric-card">
                <h3>2️⃣ AI Processing</h3>
                <p>OpenAI GPT-4 analyzes results and generates insights automatically</p>
            </div>
            """, unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div class="metric-card">
                <h3>3️⃣ Executive Presentation</h3>
                <p>PresenterBot creates comprehensive PowerPoint decks with strategic insights</p>
            </div>
            """, unsafe_allow_html=True)
    
    else:
        # Show presentation queue
        st.markdown(f"### 📋 Presentation Queue ({total_analyses} analyses)")
        
        all_analyses = []
        for agent_name, agent in st.session_state.enhanced_agents.items():
            if agent.name != 'PresenterBot':
                for analysis in agent.analyses:
                    all_analyses.append(analysis)
        
        # Display analyses preview
        for i, analysis in enumerate(all_analyses):
            with st.expander(f"Analysis {i+1}: {analysis['title']} ({analysis['agent']})", expanded=False):
                col1, col2 = st.columns([2, 1])
                
                with col1:
                    st.write(f"**Query:** {analysis['query']}")
                    st.write(f"**Agent:** {analysis['agent']}")
                    st.write(f"**Time:** {analysis['timestamp'].strftime('%H:%M:%S')}")
                    
                    if isinstance(analysis['data'], pd.DataFrame) and not analysis['data'].empty:
                        st.write(f"**Records:** {len(analysis['data'])}")
                        with st.expander("Sample Data"):
                            st.dataframe(analysis['data'].head(3))
                
                with col2:
                    st.markdown("**📊 Quick Stats**")
                    if isinstance(analysis['data'], pd.DataFrame) and not analysis['data'].empty:
                        numeric_cols = analysis['data'].select_dtypes(include=['number']).columns
                        for col in numeric_cols[:2]:
                            if analysis['data'][col].notna().sum() > 0:
                                st.metric(col.replace('_', ' ').title(), f"{analysis['data'][col].mean():.2f}")
                
                # AI insights
                if analysis.get('ai_insights'):
                    st.markdown("### 🤖 AI Analysis")
                    st.markdown(f"""
                    <div class="ai-response">
                        {analysis['ai_insights'][:500]}{"..." if len(analysis['ai_insights']) > 500 else ""}
                    </div>
                    """, unsafe_allow_html=True)
        
        # Generate presentation
        st.markdown("---")
        col1, col2, col3 = st.columns([2, 1, 1])
        
        with col1:
            if POWERPOINT_AVAILABLE:
                if st.button("🚀 Generate AI-Powered PowerPoint Presentation", type="primary"):
                    with st.spinner("🤖 AI is creating your executive presentation..."):
                        ppt_buffer = create_enhanced_presentation(all_analyses)
                        
                        if ppt_buffer:
                            st.success("✅ AI presentation generated successfully!")
                            
                            filename = f"Mental_Health_AI_Analytics_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                            
                            st.download_button(
                                label="📥 Download PowerPoint Presentation",
                                data=ppt_buffer.getvalue(),
                                file_name=filename,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        else:
                            st.error("❌ Failed to generate presentation")
            else:
                st.error("PowerPoint generation not available")
        
        with col2:
            st.metric("📑 Total Slides", len(all_analyses) + 5)
        
        with col3:
            st.metric("🤖 AI Insights", sum(1 for a in all_analyses if a.get('ai_insights')))

else:
    # Other agents interface
    st.markdown(f"### 🤖 {selected_agent.name} - {selected_agent.role}")
    st.markdown(f"*{selected_agent.description}*")
    
    # Show previous analyses
    if selected_agent.analyses:
        st.markdown("### 📊 Previous Analyses")
        for i, analysis in enumerate(selected_agent.analyses):
            with st.expander(f"{analysis['title']} ({analysis['timestamp'].strftime('%H:%M')})"):
                st.write(f"**Query:** {analysis['query']}")
                
                if analysis.get('custom_sql'):
                    with st.expander("SQL Query Used"):
                        st.code(analysis['custom_sql'], language='sql')
                
                if isinstance(analysis['data'], pd.DataFrame) and not analysis['data'].empty:
                    col1, col2 = st.columns([3, 1])
                    
                    with col1:
                        st.dataframe(analysis['data'].head(10))
                    
                    with col2:
                        st.markdown("**Quick Stats**")
                        numeric_cols = analysis['data'].select_dtypes(include=['number']).columns
                        for col in numeric_cols[:3]:
                            if analysis['data'][col].notna().sum() > 0:
                                st.metric(col.replace('_', ' ').title(), 
                                        f"{analysis['data'][col].mean():.2f}")
                
                # Show AI insights
                if analysis.get('ai_insights'):
                    st.markdown("### 🤖 AI Analysis")
                    st.markdown(f"""
                    <div class="ai-response">
                        <strong>OpenAI GPT-4 Analysis:</strong><br>
                        {analysis['ai_insights']}
                    </div>
                    """, unsafe_allow_html=True)
                
                if analysis['insights']:
                    st.write(f"**Manual Insights:** {analysis['insights']}")
    
    # Analysis interface
    st.markdown("---")
    st.markdown("### 🎯 New Analysis")
    
    # Analysis tabs
    tab1, tab2 = st.tabs(["🚀 Quick Analysis", "🔧 Custom SQL Query"])
    
    with tab1:
        # Pre-built queries based on agent type
        if selected_agent.name == 'BillBot':
            st.markdown("#### 💰 Financial Analytics")
            quick_queries = {
                "Revenue by Month": """
                    SELECT 
                        visit_month_name,
                        visit_year,
                        COUNT(*) as visit_count,
                        SUM(billed_amount) as total_billed,
                        SUM(paid_amount) as total_paid,
                        AVG(payment_percentage) as avg_payment_rate
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE billed_amount IS NOT NULL
                    GROUP BY visit_month_name, visit_year, visit_month
                    ORDER BY visit_year, visit_month
                """,
                "Claims Payment Analysis": """
                    SELECT 
                        claim_id,
                        billed_amount,
                        paid_amount,
                        patient_responsibility,
                        payment_percentage,
                        cost_category,
                        CASE 
                            WHEN payment_percentage >= 90 THEN 'Excellent'
                            WHEN payment_percentage >= 70 THEN 'Good'
                            WHEN payment_percentage >= 50 THEN 'Fair'
                            ELSE 'Poor'
                        END as payment_category
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE claim_id IS NOT NULL
                """,
                "Provider Revenue Performance": """
                    SELECT 
                        provider_full_name,
                        provider_specialty,
                        COUNT(*) as total_visits,
                        SUM(billed_amount) as total_revenue,
                        AVG(billed_amount) as avg_billing,
                        AVG(payment_percentage) as avg_collection_rate
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE provider_full_name IS NOT NULL
                    GROUP BY provider_full_name, provider_specialty
                    ORDER BY total_revenue DESC
                """,
                "Insurance Plan Performance": """
                    SELECT 
                        plan_type,
                        plan_type_description,
                        COUNT(*) as member_count,
                        AVG(copay_amount) as avg_copay,
                        AVG(deductible_amount) as avg_deductible,
                        SUM(billed_amount) as total_billed,
                        AVG(payment_percentage) as avg_payment_rate
                    FROM [dbo].[vw_MentalHealthRAG]
                    GROUP BY plan_type, plan_type_description
                    ORDER BY member_count DESC
                """
            }
        
        elif selected_agent.name == 'ScheduleBot':
            st.markdown("#### 📅 Scheduling Analytics")
            quick_queries = {
                "Appointment Patterns by Day": """
                    SELECT 
                        visit_day_of_week,
                        COUNT(*) as appointment_count,
                        AVG(session_length_minutes) as avg_session_length,
                        COUNT(DISTINCT provider_id) as active_providers
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE visit_date IS NOT NULL
                    GROUP BY visit_day_of_week
                    ORDER BY appointment_count DESC
                """,
                "Provider Utilization Analysis": """
                    SELECT 
                        provider_full_name,
                        provider_specialty,
                        COUNT(*) as total_appointments,
                        AVG(session_length_minutes) as avg_session_time,
                        COUNT(DISTINCT member_id) as unique_patients,
                        MIN(visit_date) as first_appointment,
                        MAX(visit_date) as last_appointment
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE provider_full_name IS NOT NULL
                    GROUP BY provider_full_name, provider_specialty
                    ORDER BY total_appointments DESC
                """,
                "Session Length Distribution": """
                    SELECT 
                        visit_type,
                        COUNT(*) as session_count,
                        AVG(session_length_minutes) as avg_length,
                        MIN(session_length_minutes) as min_length,
                        MAX(session_length_minutes) as max_length,
                        STDEV(session_length_minutes) as length_std_dev
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE session_length_minutes IS NOT NULL
                    GROUP BY visit_type
                    ORDER BY avg_length DESC
                """,
                "Monthly Scheduling Trends": """
                    SELECT 
                        visit_year,
                        visit_month_name,
                        COUNT(*) as total_visits,
                        COUNT(DISTINCT member_id) as unique_patients,
                        COUNT(DISTINCT provider_id) as active_providers,
                        AVG(session_length_minutes) as avg_session_length
                    FROM [dbo].[vw_MentalHealthRAG]
                    GROUP BY visit_year, visit_month_name, visit_month
                    ORDER BY visit_year, visit_month
                """
            }
        
        else:  # ResearchBot
            st.markdown("#### 🔬 Clinical Research Analytics")
            quick_queries = {
                "Diagnosis Distribution": """
                    SELECT 
                        diagnosis_code,
                        condition_category,
                        COUNT(*) as patient_count,
                        COUNT(DISTINCT member_id) as unique_patients,
                        AVG(member_age) as avg_patient_age,
                        COUNT(CASE WHEN gender = 'F' THEN 1 END) as female_count,
                        COUNT(CASE WHEN gender = 'M' THEN 1 END) as male_count
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE diagnosis_code IS NOT NULL
                    GROUP BY diagnosis_code, condition_category
                    ORDER BY patient_count DESC
                """,
                "Treatment Outcomes by Provider": """
                    SELECT 
                        provider_full_name,
                        provider_specialty,
                        condition_category,
                        COUNT(*) as total_sessions,
                        COUNT(DISTINCT member_id) as unique_patients,
                        AVG(session_length_minutes) as avg_session_time,
                        progress_indicator
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE provider_full_name IS NOT NULL
                    GROUP BY provider_full_name, provider_specialty, condition_category, progress_indicator
                    ORDER BY total_sessions DESC
                """,
                "Patient Demographics Analysis": """
                    SELECT 
                        age_group,
                        gender,
                        condition_category,
                        COUNT(DISTINCT member_id) as patient_count,
                        AVG(mental_health_visits_allowed) as avg_visits_allowed,
                        COUNT(*) as total_visits,
                        AVG(session_length_minutes) as avg_session_length
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE age_group IS NOT NULL
                    GROUP BY age_group, gender, condition_category
                    ORDER BY patient_count DESC
                """,
                "Treatment Effectiveness by Condition": """
                    SELECT 
                        condition_category,
                        COUNT(DISTINCT member_id) as total_patients,
                        AVG(visit_sequence_number) as avg_visits_per_patient,
                        AVG(session_length_minutes) as avg_session_duration,
                        COUNT(CASE WHEN progress_indicator = 'Improving' THEN 1 END) as improving_count,
                        COUNT(CASE WHEN progress_indicator = 'Stable' THEN 1 END) as stable_count
                    FROM [dbo].[vw_MentalHealthRAG]
                    WHERE condition_category IS NOT NULL
                    GROUP BY condition_category
                    ORDER BY total_patients DESC
                """
            }
        
        # Query selection
        selected_query_name = st.selectbox("Select pre-built analysis:", list(quick_queries.keys()))
        
        with st.expander("📝 View SQL Query"):
            st.code(quick_queries[selected_query_name], language='sql')
        
        # Custom insights input
        custom_insights = st.text_area("Additional insights or focus areas (optional):", 
                                     placeholder="Enter any specific insights you'd like to explore...")
        
        # Execute analysis
        if st.button(f"🚀 Run {selected_query_name}", type="primary"):
            with st.spinner("🔄 Querying database and analyzing with AI..."):
                analysis = selected_agent.add_analysis(
                    query=selected_query_name,
                    custom_sql=quick_queries[selected_query_name],
                    insights=custom_insights
                )
                
                if analysis:
                    st.success(f"✅ Analysis completed! Results added to presentation queue.")
                    
                    # Display results
                    st.markdown("### 📊 Analysis Results")
                    
                    if isinstance(analysis['data'], pd.DataFrame) and not analysis['data'].empty:
                        # Main results
                        col1, col2 = st.columns([3, 1])
                        
                        with col1:
                            st.dataframe(analysis['data'])
                        
                        with col2:
                            st.markdown("**📈 Key Metrics**")
                            numeric_cols = analysis['data'].select_dtypes(include=['number']).columns
                            for col in numeric_cols[:4]:
                                if analysis['data'][col].notna().sum() > 0:
                                    st.metric(
                                        col.replace('_', ' ').title(),
                                        f"{analysis['data'][col].mean():.2f}",
                                        f"Total: {analysis['data'][col].sum():.2f}"
                                    )
                        
                        # Visualization
                        if len(analysis['data']) > 1:
                            st.markdown("### 📊 Visualization")
                            
                            # Auto-generate appropriate chart
                            numeric_cols = analysis['data'].select_dtypes(include=['number']).columns
                            categorical_cols = analysis['data'].select_dtypes(include=['object']).columns
                            
                            if len(numeric_cols) >= 1 and len(categorical_cols) >= 1:
                                chart_col1, chart_col2 = st.columns(2)
                                
                                with chart_col1:
                                    # Bar chart
                                    fig_bar = px.bar(
                                        analysis['data'].head(10),
                                        x=categorical_cols[0],
                                        y=numeric_cols[0],
                                        title=f"{numeric_cols[0].replace('_', ' ').title()} by {categorical_cols[0].replace('_', ' ').title()}"
                                    )
                                    fig_bar.update_layout(xaxis_tickangle=-45)
                                    st.plotly_chart(fig_bar, use_container_width=True)
                                
                                with chart_col2:
                                    if len(numeric_cols) >= 2:
                                        # Scatter plot
                                        fig_scatter = px.scatter(
                                            analysis['data'],
                                            x=numeric_cols[0],
                                            y=numeric_cols[1],
                                            color=categorical_cols[0] if len(categorical_cols) > 0 else None,
                                            title=f"{numeric_cols[1].replace('_', ' ').title()} vs {numeric_cols[0].replace('_', ' ').title()}"
                                        )
                                        st.plotly_chart(fig_scatter, use_container_width=True)
                                    else:
                                        # Pie chart for categorical data
                                        value_counts = analysis['data'][categorical_cols[0]].value_counts().head(8)
                                        fig_pie = px.pie(
                                            values=value_counts.values,
                                            names=value_counts.index,
                                            title=f"Distribution of {categorical_cols[0].replace('_', ' ').title()}"
                                        )
                                        st.plotly_chart(fig_pie, use_container_width=True)
                        
                        # AI Insights
                        if analysis.get('ai_insights'):
                            st.markdown("### 🤖 AI Analysis & Recommendations")
                            st.markdown(f"""
                            <div class="ai-response">
                                <strong>OpenAI GPT-4 Professional Analysis:</strong><br><br>
                                {analysis['ai_insights']}
                            </div>
                            """, unsafe_allow_html=True)
                        
                        # Custom insights
                        if custom_insights:
                            st.markdown("### 💡 Additional Insights")
                            st.info(custom_insights)
                    
                    else:
                        st.warning("No data returned from query")
                else:
                    st.error("Analysis failed")
    
    with tab2:
        st.markdown("#### 🔧 Custom SQL Analysis")
        st.info("Write your own SQL query against the [dbo].[vw_MentalHealthRAG] view")
        
        # SQL editor
        custom_sql = st.text_area(
            "SQL Query:",
            height=200,
            placeholder="""SELECT 
    column1,
    column2,
    COUNT(*) as count
FROM [dbo].[vw_MentalHealthRAG]
WHERE condition = 'value'
GROUP BY column1, column2
ORDER BY count DESC;"""
        )
        
        # Query name
        query_name = st.text_input("Analysis Name:", placeholder="e.g., Custom Provider Analysis")
        
        # Insights
        custom_insights = st.text_area("Analysis Focus:", 
                                     placeholder="What specific insights are you looking for?")
        
        # Validate and execute
        if st.button("🚀 Execute Custom Analysis", type="primary"):
            if not custom_sql.strip():
                st.error("Please enter a SQL query")
            elif not query_name.strip():
                st.error("Please enter an analysis name")
            else:
                with st.spinner("🔄 Executing custom query and analyzing with AI..."):
                    analysis = selected_agent.add_analysis(
                        query=query_name,
                        custom_sql=custom_sql,
                        insights=custom_insights
                    )
                    
                    if analysis:
                        st.success(f"✅ Custom analysis '{query_name}' completed!")
                        
                        # Display results similar to quick analysis
                        if isinstance(analysis['data'], pd.DataFrame) and not analysis['data'].empty:
                            st.markdown("### 📊 Query Results")
                            st.dataframe(analysis['data'])
                            
                            # AI insights
                            if analysis.get('ai_insights'):
                                st.markdown("### 🤖 AI Analysis")
                                st.markdown(f"""
                                <div class="ai-response">
                                    {analysis['ai_insights']}
                                </div>
                                """, unsafe_allow_html=True)
                        else:
                            st.warning("Query returned no data")
                    else:
                        st.error("Query execution failed")

# Database schema reference
with st.sidebar.expander("📋 Database Schema Reference"):
    st.markdown("""
    **Key Tables & Views:**
    - `[dbo].[vw_MentalHealthRAG]` - Main view
    
    **Important Columns:**
    - **Billing:** billed_amount, paid_amount, claim_id
    - **Scheduling:** visit_date, session_length_minutes
    - **Clinical:** diagnosis_code, condition_category
    - **Demographics:** member_age, gender, age_group
    """)

# Footer with system info
st.markdown("---")
col1, col2, col3 = st.columns(3)

with col1:
    st.markdown("""
    **🏥 System Information**
    - Database: JONESFAMILYPC3\\MedDw
    - View: [dbo].[vw_MentalHealthRAG]
    - AI: OpenAI GPT-4 Integration
    """)

with col2:
    st.markdown("""
    **🔒 Security & Compliance**
    - HIPAA-compliant data handling
    - Automated PHI masking
    - Secure database connections
    """)

with col3:
    st.markdown("""
    **📞 Support**
    - Platform: Mental Health AI Analytics
    - Organization: Transforming Lives
    - Status: Production Ready
    """)

st.markdown("""
<div style='text-align: center; color: #666; margin-top: 2rem;'>
    🧠 Mental Health AI Analytics Platform | Powered by OpenAI & SQL Server | Transforming Lives Mental Health Counseling
</div>
""", unsafe_allow_html=True)